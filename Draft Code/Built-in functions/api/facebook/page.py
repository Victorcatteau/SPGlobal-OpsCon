#!/usr/bin/env python
# -*- coding: utf-8 -*-

'''---------------------------------------------------------------------
    FACEBOOK SCRAPPING
    author : Theo ALVES DA COSTA
    date started : 01/12/2016 
   ---------------------------------------------------------------------
'''


#-------------------------------------------------------------------------
# LIBRARIES

# Usual
import numpy as np
import pandas as pd
import time
import os
import matplotlib.pyplot as plt
import requests
import json
import datetime
from tqdm import tqdm
from urllib import request
from collections import defaultdict
from IPython.display import display,HTML,Markdown

# Others
import unidecode
import bs4
from scipy.interpolate import UnivariateSpline

# Ekimetrics library
from ekimetrics.nlp.utils import extract_emojis
from ekimetrics.nlp.models import *
from ekimetrics.utils.time import sec_to_hms
from ekimetrics.utils import gender
from ekimetrics.visualizations import charts

from ekimetrics.api.facebook.posts import *
from ekimetrics.api.facebook.users import *
from ekimetrics.api.facebook.comments import *
from ekimetrics.api.facebook.utils import *


try:
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION,XL_LABEL_POSITION
    from pptx.util import Inches
except:
    print("Warning, skipping import of pptx - install it to enjoy ppt capabilities")


#================================================================================================================================================
# PAGE
#================================================================================================================================================

class Page(object):
    """
    Page data representation
    """

    def __init__(self, page_id="nissanelectric", n_posts=None, details=True, file_path=None,token = None):
        """
        Initialization
        """

        s = time.time()

        # RELOADING THE DATA VIA FACEBOOK API
        if file_path is None:

            # Connection to the page data
            self.details = details
            self.page_id = page_id
            page_data = self.build_connection(token = token)

            # Retrieving the basic info
            self.name = page_data.get('name')
            self.about = page_data.get('about')
            self.followers_count = page_data.get("fan_count")

            # Getting all the posts
            self._print('##### Retrieving posts on the feed of ``{}`` facebook page'.format(self.name))
            self.connection_done = False
            self.posts = self.get_all_posts(n_posts)

            # Parsing all the posts
            self._print('##### Parsing every posts data')
            for i, post in enumerate(self.posts):
                try:
                    self.posts[i] = Post(post, details=details, page=self.name, post_number=i, posts_count=len(self.posts))
                except Exception as e:
                    print("- ERROR - ",e)

            # Protection in case all the posts were not retrieved
            self.posts = [post for post in self.posts if type(post) != dict]

            print('\n... Done')

            # Saving automatically the data for further use
            self.save_data()

            # Finishing the process
            e = time.time()
            self._print('##### Scrapping done in ``{}``'.format(sec_to_hms(e - s)))

        # LOADING THE DATA FROM A JSON FILE
        else:
            data = self.load_data(file_path)
            self.page_id = data["page_id"]
            self.details = data["details"]
            self.name = data["name"]
            self.about = data["about"]
            self.posts = [Post(post, reload=True, page=self.name)for post in data["posts"]]

        # Other parameters
        self.genders_detected = False
        self.ages_detected = False




    #-------------------------------------------------------------------------
    # OPERATORS

    def __repr__(self):
        """
        Representation operator
        """
        return self.about

    def __str__(self):
        """
        String representation operator
        """
        return self.about


    def __iter__(self):
        """
        Post iterator operator
        """
        return iter(self.posts)


    def __getitem__(self,key):
        """
        Get a post
        """
        return self.posts[key]


    def __len__(self):
        """
        Return number of posts
        """
        return len(self.posts)


    def _print(self,message):
        """
        Print function with markdown in notebook
        """
        if self._is_interactive():
            display(Markdown(message))
        else:
            print(message)

    def _is_interactive(self):
        """
        Boolean function that checks if we are running the score in a notebook or in a python script
        """
        import __main__ as main
        return not hasattr(main, '__file__')



    #-------------------------------------------------------------------------
    # IO

    def save_data(self):
        """
        Save all the data as a json file
        """
        file_name = self.page_id + "_" + str(pd.to_datetime("today"))[:10] + ".json"
        with open(file_name, 'w') as file:
            json.dump(self.get_data(), file)




    def load_data(self, file_path):
        """
        Reload data saved in a json file
        """
        return json.loads(open(file_path, "r").read())






    def get_data(self):
        """
        Export the data of all posts in a page as a list of dictionary by posts
        Is used to save the data in a json file
        :returns: list of dicts -- with all the page data
        """

        data = {"page_id": self.page_id, "name": self.name,"about": self.about, "details": self.details}
        data["posts"] = [post.get_data() for post in self.posts]
        return data




    def export_excel_files(self):
        pass



    #-------------------------------------------------------------------------
    # CONSTRUCTION

    def build_connection(self,token = None):
        """
        Create the connection parameters required to connecte to the facebook API
        :returns: dict -- a json dictionary with the API output 
        """
        self.connection = FacebookConnection(token = token)
        self.access_token = self.connection.access_token
        self.base_url = "https://graph.facebook.com/v2.8"
        self.node_url_feed = "/{}/feed".format(self.page_id)
        self.node_url_posts = "/{}/posts".format(self.page_id)
        url = "{}/{}?fields=about,name,fan_count&access_token={}".format(self.base_url, self.page_id, self.access_token)
        page_data = self.connection.get_data(url)
        return page_data



    def build_url(self, n_posts):
        n_details = 100 if self.details else 1
        reactions = "limit({0}).summary(true)".format(n_details)
        comments = "limit({0}).summary(true){{created_time,from,message,id,message_tags}}".format(n_details)
        parameters = ("/?fields=message,link,created_time,type,from,name,id,full_picture,"
                      "likes.limit({0}).summary(true),reactions.{3},"
                      "comments.limit({0}).summary(true){{created_time,from,message,id,message_tags,"
                      "reactions.{3},comments.{4}}},"
                      "shares&limit={1}&access_token={2}"
                      .format(n_details, n_posts, self.access_token, reactions, comments))
        return self.base_url + self.node_url_posts + parameters





    #-------------------------------------------------------------------------
    # GETTERS

    def get_n_posts(self, n_posts=100, url=None):
        # n = 100 at max
        n_posts = min(100, n_posts)

        if url is None:
            url = self.build_url(n_posts)

        data = self.connection.get_data(url, self.connection_done)
        while not self.connection_done:
            if data:
                self.connection_done = True
            else:
                print(
                    '... Impossible to connect with {} posts. Trying with 5 less.'.format(n_posts))
                n_posts -= 5
                url = self.build_url(n_posts)
                data = self.connection.get_data(url, self.connection_done)

        return data











    def get_all_posts(self, n_posts):
        posts = []
        left_posts = n_posts if n_posts is not None else np.inf
        url = None
        scrapping = True

        while True:
            if left_posts < 100:
                data = self.get_n_posts(n_posts, url)
                posts_processed = len(data['data'])
                posts += data['data']
                print('... {} posts retrieved'.format(len(posts)), end='\r')
                break
            else:
                # Retrieve 25 posts
                data = self.get_n_posts(25, url)
                posts_processed = len(data['data'])
                if posts_processed:
                    posts += data['data']
                else:
                    break

                left_posts = left_posts - 25
                print('... {} posts retrieved'.format(len(posts)), end='\r')

                try:
                    url = data['paging']['next']
                except KeyError:
                    break

        print("... Total : {} posts retrieved. Scrapping done.".format(len(posts)))

        return posts







    def get_all_admin_posts(self, date_start = None,date_end = None):
        """
        Returns all the admin posts made between 2 dates 
        Returns every posts if no dates are provided
        Uses the operator surchage that compares the date of a post with a date
        :param str date_start: starting date
        :param str date_end: end date
        :returns: list of Post objects
        """
        return [p for p in self.posts if (p.admin and p > date_start and p < date_end)]






    def get_all_matchReactionComment(self):
        tousLesMatchs = {'SAD': [], 'LIKE': [],
                         'LOVE': [], 'ANGRY': [], 'WOW': [], 'HAHA': []}
        for post in self.posts:
            lesMatchsduPost = post.get_intersection_reactions_comments()
            for typeReaction in lesMatchsduPost.keys():
                tousLesMatchs[typeReaction].append(
                    lesMatchsduPost[typeReaction])
        return(tousLesMatchs)






    def get_all_reactions(self):
        all_reaction = {'SAD':[],'LIKE':[],
                        'LOVE':[],'ANGRY':[],'WOW':[], 'HAHA':[]}
        for post in self.posts :
            for reaction in post.reactions.keys():
                all_reaction[reaction].append(post.reactions[reaction])
        return(all_reaction)









    #-------------------------------------------------------------------------
    # ANALYSIS






    def compute_performances(self,monthly=True):
        # GET DATA
        data = [{"date": post.date, **post.get_performances()} for post in self.posts if post.admin]
        reactions = ["likes", "comments", "love","haha", "wow", "angry", "sad"]
        output = pd.DataFrame(data, columns=['date'] + reactions)

        # MONTHLY ANALYSIS
        if monthly:
            output["date"] = list(map(lambda x: pd.to_datetime("{}/{}/01".format(x.year, x.month)),pd.to_datetime(output['date'])))
        else:
            output["date"] = list(map(lambda x: pd.to_datetime("{}/{}/{}".format(x.year, x.month, x.day)),pd.to_datetime(output['date'])))

        # GROUPBY
        output = (output.groupby('date', as_index=False)
                  .agg({
                      'likes': {'count': 'count', 'likes': 'sum'},
                      'comments': {'comments': 'sum'},
                      'love': {'love': 'sum'},
                      'haha': {'haha': 'sum'},
                      'wow': {'wow': 'sum'},
                      'angry': {'angry': 'sum'},
                      'sad': {'sad': 'sum'}})
                  .sort_values('date', ascending=True))
        output.columns = output.columns.get_level_values(1)
        output.columns = ['date'] + list(output.columns)[1:]
        output = output.loc[output.date > "2010/01/01"]
        output.fillna(0, inplace=True)
        output.set_index("date", inplace=True)

        # CREATE NEW KPIs
        for reaction in reactions:
            output['average_{}'.format(reaction)] = output[reaction] / output['count'].astype(float)

        average_reactions = ["average_" + reaction for reaction in reactions]


        return output



    def show_performances(self,monthly=True,with_plotly=True, on_notebook=True,as_average_per_post = True,**kwargs):

        performances = self.compute_performances(monthly = monthly)
        reactions = ["likes", "comments", "love","haha", "wow", "angry", "sad"]
        average_reactions = ["average_" + reaction for reaction in reactions]

        if as_average_per_post:
            performances = performances[average_reactions]
            performances.columns = reactions
        else:
            performances = performances[reactions]


        # PLOT
        if plot:
            fig = charts.plot_line_chart(performances, with_plotly=with_plotly, on_notebook=on_notebook,
                                            title = "{}'s page posts performances".format(self.name),
                                            xlabel = "Date",ylabel = "Number",
                                            **kwargs)
            return fig

        # RETURN RESULTS
        else:
            return output






    def find_top_posts(self, axis = "likes", n=None):
        """
        Find the best posts along a indicator axis as a dataframe holding all the necessary information
        """

        # Define fields list
        fields = ["id", "date", "creator", "type", "message","link", "picture"]
        reactions = ["likes", "comments", "love","haha", "wow", "angry", "sad"]
        assert axis in fields + reactions

        # Helper function
        def get_data(post):
            d = {k: v for k, v in post.__dict__.items() if k in fields}
            d = {**d,**post.get_performances()}
            return d

        # Extract the KPIs for all posts
        dictionary_of_posts = [get_data(post) for post in self.posts]

        # Convert the dictionary 
        output_df = pd.DataFrame(dictionary_of_posts, columns=fields+reactions).fillna(0.0)

        # Sort along the axis
        return output_df.sort_values(axis, ascending=False).head(n)





    def compute_posts_type(self,rename = False):
        """
        
        """
        types = [post.type for post in self.posts]
        df = pd.DataFrame(types,columns = ["type"])
        df["count"] = 1
        df = df.groupby("type").sum().sort_values("count",ascending = False)
        if rename:
            df = df.rename(columns = {"count":self.name})
        return df



    def compute_followers_engagement(self):
        
        # Compute fanbase
        if not hasattr(self,"fanbase"):
            self.build_users_fanbase()

        # Return engagement
        return  len(self.fanbase.data) / self.followers_count




    def show_posts_type(self,kind = "pie",with_plotly = True,on_notebook = True,**kwargs):
        assert kind in ["pie","bar"]

        df = self.compute_posts_type()

        if kind == "pie":
            fig = charts.plot_pie_chart(df,with_plotly = with_plotly,on_notebook = on_notebook,
                                        title = "{} post type distribution".format(self.name),
                                        **kwargs)
        elif kind == "bar":
            fig = charts.plot_bar_chart(df,with_plotly = with_plotly,on_notebook = on_notebook,
                                        title = "{} post type distribution".format(self.name),
                                        xlabel = "Type",ylabel = "Number of posts",
                                        **kwargs)

        return fig








    #-------------------------------------------------------------------------------------------
    # CRM ANALYSIS


    def build_users_fanbase(self,date_start=None,date_end = None,limit=None,inplace = True,verbose = 1):
        """
        Get aggregated data of users interacting with the page
        Can be used to find the core fan base, or study the influencers for example
        """

        if date_start is None and date_end is None and hasattr(self,"fanbase"):
            if not inplace:
                return self.fanbase
        else:
            if verbose: print(">> Calculating {}'s fanbase {}".format(self.name,Period(date_start,date_end).__repr__()))

            # Get all the posts
            posts = self.get_all_admin_posts(date_start = date_start, date_end = date_end)

            # Define a user dictionary
            users = defaultdict(dict)

            # Iterate on the posts
            for post in posts:

                # Iterate on the type of reaction
                for reaction in post.reactions:

                    # Find all the users
                    for user in post.reactions[reaction]:
                        users[user.id].setdefault('user', user)
                        users[user.id][reaction] = users[user.id].setdefault(reaction, 0) + 1

                # Iterate on the comments too
                for comment in post.comments:
                    users[comment.id].setdefault('user', comment.user)
                    users[comment.id]['comments'] = users[comment.id].setdefault('comments', 0) + 1

            # Define lists of interactions
            all_reactions = ['LIKE', 'HAHA', 'LOVE', 'SAD', 'WOW', 'ANGRY']
            all_interactions = ['LIKE', 'HAHA', 'LOVE', 'SAD', 'WOW', 'ANGRY', 'comments']
            all_columns = ["user"] + all_interactions

            # Create a pandas DataFrame
            df = pd.DataFrame(users).transpose().fillna(0)
            columns = [col for col in all_columns if col in df.columns]
            df = df[columns]

            # Get only existing columns
            all_reactions = [x for x in all_reactions if x in columns]
            all_interactions = [x for x in all_interactions if x in columns]

            # Calculate the sum of total reactions and interactions
            df["total_reactions"] = df[all_reactions].sum(axis = 1)
            df["total_interactions"] = df[all_interactions].sum(axis = 1)
            df["reactions_ratio"] = (df["total_reactions"] / len(posts)).map(lambda x:round(x,3))
            df.sort_values('total_interactions', ascending=False, inplace=True)

            # Return the result
            data = df.head(limit)
            fanbase = Fanbase(data = data,page_name = self.name)

            if inplace:
                self.fanbase = fanbase
            else:
                return fanbase




    def compute_fanbase_engagement(self,top = 0.1,periods = [Period()],axis = "total_reactions"):
        """
        Find the gender engagement
        """
        print(">> Getting fanbase engagement")

        # Assertions
        assert isinstance(periods,list)


        # Data initialization
        engagement = defaultdict(dict)

        # Iteration over the periods
        for period in periods:

            # Create a key for the period
            period_str = str(period)

            # Create the fanbase
            fanbase = self.build_users_fanbase(**period.get(),inplace = False)

            # Get the engagement for the given fanbase
            engagement[period_str] = fanbase.compute_engagement()


        engagement = pd.DataFrame(engagement).transpose()
        engagement = engagement[[f for f in engagement.columns if f.startswith("top")]+["100%"]]


        return engagement





    def show_fanbase_engagement(self,engagement,with_plotly = True,on_notebook = True,normalized = True,kind = "area",**kwargs):
        assert kind in ["area","bar"]

        # Plot bar chart
        if kind == "bar" or len(engagement) == 1:
            fig = charts.plot_bar_chart(engagement,with_plotly = with_plotly,on_notebook = on_notebook)

        # Plot filled areas chart
        else:
            ylabel = "Number of interactions"
            engagement = engagement.copy()
            if normalized:
                for column in engagement.columns:
                    engagement[column] = engagement[column] / engagement["100%"]
                ylabel += " (%)"
            fig = charts.plot_area_chart(engagement,with_plotly = with_plotly,on_notebook = on_notebook,
                                            title = "{}'s users engagement".format(self.name),
                                            xlabel = "Period",ylabel = ylabel,
                                            **kwargs)


        return fig





    #-------------------------------------------------------------------------
    # NLP ANALYSIS


    def build_posts_corpus(self,**kwargs):
        """
        Build a corpus ontology from the message posts published by the page 
        """
        texts = [post.message for post in self.posts]
        return Corpus(texts = texts,**kwargs)



    def build_comments_corpus(self):
        """
        Build a corpus ontology from the comments published by the users
        """
        comments = [comment for document in self.posts for comment in document.comments]
        return Comments(comments=comments)






    #-------------------------------------------------------------------------
    # POWERPOINT ANALYSIS



    def create_presentation(self):
        """
        Start building the Ekimetrics presentation
        """
        from ekimetrics.utils.slider import EkimetricsSlider
        self.slider = EkimetricsSlider()



    def save_presentation(self):
        """
        Save the presentation as a Powerpoint presentation
        """
        file_name = self.page_id + "_" + str(pd.to_datetime("today"))[:10] + ".pptx"
        self.slider.save(file_name)



    def build_presentation(self):
        """
        Build the final presentation with all the analysis
        """
        self.create_presentation()
        self.add_slide_first_slide()
        self.add_slide_introduction()
        self.add_slide_posts_type()




    def add_main_slide(self,title):
        """
        Helper function to build a blank slide
        """
        slide = self.slider.create_base_main_slide(title,"{} Facebook Analysis".format(self.name))
        return slide


    def add_slide_first_slide(self):
        """
        Slide 1 : Title
        """
        self.slider.create_base_title_slide("{} Facebook Analysis\n{}".format(self.name,str(pd.to_datetime("today"))[:10]))



    def add_slide_introduction(self):
        slide = self.add_main_slide("Introduction")
        



    def add_slide_posts_type(self):
        """
        Slide : Posts by type
        """
        slide = self.add_main_slide("Analysis by post type")

        df = self.compute_posts_type()


        chart_data = ChartData()
        chart_data.categories = df.index
        chart_data.add_series('Count', tuple(df["count"]))

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        # data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END





    #-------------------------------------------------------------------------
    # COMPUTER VISION


    def download_all_pictures(self,base_path = None,store_as_attribute = False):

        # Create base path if not given
        if base_path is None:
            base_path = "image_{}".format(self.name.lower())

        if not os.path.exists(base_path):
            os.mkdir(base_path)


        # Iterate over the posts
        for post in tqdm(self.posts,desc = "Downloading all pictures"):
            try:
                post.download_picture(base_path = base_path,path = None,store_as_attribute = store_as_attribute)
            except Exception as e:
                print(e)


    def store_all_pictures(self):
        """
        Store each post picture as attributes of the post objects
        """

        # Iterate over the posts
        for post in tqdm(self.posts,desc = "Storing all pictures"):
            try:
                post.store_picture()
            except Exception as e:
                post.img = None
        



    def analyze_all_pictures(self,model = None):

        # Retrieving the model
        if model is None:
            model = ComputerVisionModel(name = "vgg16")

        else:
            if type(model) == str:
                model = ComputerVisionModel(name = model)


        # Make predictions
        predictions = {}

        for post in tqdm(self.posts,desc = "Analyzing all pictures"):
            predictions[self.id] = post.analyze_picture(model = model)


        return predictions



