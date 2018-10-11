#!/usr/bin/env python
# -*- coding: utf-8 -*- 


"""--------------------------------------------------------------------
PYTHON POWERPOINT SLIDER

Started on the 20/10/2017

Creators : 
    - Theo ALVES DA COSTA

------------------------------------------------------------------------
"""


from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


#=============================================================================================================================
# Ekimetrics Slider Wrapper
#=============================================================================================================================






class EkimetricsSlider(object):
    """
    You can create or update ppt presentation using EkimetricsSlider class

    """
    def __init__(self, file_path = "C:/data/slider/ekimetrics_template.pptx"):
        """
        file_path could be either an existing presentation that you want to update OR a template for your new presentation 
        """
        self.presentation = Presentation(file_path)


    def save(self,file_path):
        self.presentation.save(file_path)



    def print_known_placeholders(self,slide):
        for shape in slide.placeholders:
            print('ID : %d - %s' % (shape.placeholder_format.idx, shape.name))




    def create_base_title_slide(self,title = "None"):
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        slide_title = slide.placeholders[10]
        slide_title.text = title
        return slide


    def create_base_section_slide(self):
        pass


        
    def create_base_main_slide(self,title="None",subtitle = "0. None"):
        title_slide_layout = self.presentation.slide_layouts[3]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        slide_title = slide.shapes.title
        slide_subtitle = slide.placeholders[13]
        slide_title.text = title_slide_layout
        slide_subtitle.text = subtitle
        return slide


    def add_picture(self, img_path, slide_number, left = Inches(2), top = Inches(2), height = Inches(1)):
        
        slide = self.presentation.slides[slide_number]
        pic = slide.shapes.add_picture(img_path, left, top, height= height)

        return slide




